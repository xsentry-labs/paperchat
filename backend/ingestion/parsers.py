"""Document parsers for all supported file types."""
from __future__ import annotations
from dataclasses import dataclass, field


@dataclass
class PageInfo:
    text: str
    page: int


@dataclass
class ParsedDocument:
    text: str
    pages: list[PageInfo] = field(default_factory=list)


async def parse_document(content: bytes, mime_type: str, filename: str) -> ParsedDocument:
    if mime_type == "application/pdf":
        return await _parse_pdf(content)
    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _parse_docx(content)
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return _parse_pptx(content)
    elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return _parse_xlsx(content)
    elif mime_type in ("text/html", "application/xhtml+xml"):
        return _parse_html(content)
    elif mime_type == "application/epub+zip":
        return _parse_epub(content)
    else:
        # Plain text, markdown, etc.
        return _parse_text(content)


async def _parse_pdf(content: bytes) -> ParsedDocument:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=content, filetype="pdf")
    pages: list[PageInfo] = []
    full_text_parts: list[str] = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        pages.append(PageInfo(text=text, page=page_num + 1))
        full_text_parts.append(text)

    full_text = "\n\n".join(full_text_parts)

    # Detect scanned PDF (no extractable text) — not supported, return what we have
    avg_chars = len(full_text) / max(len(doc), 1)
    if avg_chars < 80:
        raise ValueError(
            "This PDF appears to be scanned (image-only) and cannot be processed. "
            "Please upload a text-based PDF or a Word document."
        )

    return ParsedDocument(text=full_text, pages=pages)


def _parse_docx(content: bytes) -> ParsedDocument:
    import io
    from docx import Document

    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs)
    return ParsedDocument(text=text)


def _parse_pptx(content: bytes) -> ParsedDocument:
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    pages: list[PageInfo] = []
    parts: list[str] = []

    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        slide_text = "\n".join(slide_texts)
        pages.append(PageInfo(text=slide_text, page=i + 1))
        parts.append(slide_text)

    return ParsedDocument(text="\n\n".join(parts), pages=pages)


def _parse_xlsx(content: bytes) -> ParsedDocument:
    import io
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(",".join(cells))
        if rows:
            parts.append(f"Sheet: {sheet_name}\n" + "\n".join(rows))

    return ParsedDocument(text="\n\n".join(parts))


def _parse_html(content: bytes) -> ParsedDocument:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(content, "lxml")
    # Remove script/style tags
    for tag in soup(["script", "style", "meta", "link"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    # Clean up blank lines
    lines = [l.strip() for l in text.splitlines()]
    text = "\n".join(l for l in lines if l)
    return ParsedDocument(text=text)


def _parse_epub(content: bytes) -> ParsedDocument:
    import io
    import zipfile
    from xml.etree import ElementTree as ET

    parts: list[str] = []

    with zipfile.ZipFile(io.BytesIO(content)) as zf:
        # Find OPF file
        container = zf.read("META-INF/container.xml")
        root = ET.fromstring(container)
        ns = {"c": "urn:oasis:names:tc:opendocument:xmlns:container"}
        opf_path = root.find(".//c:rootfile", ns).get("full-path")

        opf_dir = "/".join(opf_path.split("/")[:-1])
        opf = ET.fromstring(zf.read(opf_path))
        opf_ns = {"opf": "http://www.idpf.org/2007/opf"}

        # Build manifest
        manifest: dict[str, str] = {}
        for item in opf.findall(".//opf:manifest/opf:item", opf_ns):
            manifest[item.get("id")] = item.get("href")

        # Follow spine order
        spine = opf.find(".//opf:spine", opf_ns)
        for itemref in spine.findall("opf:itemref", opf_ns):
            idref = itemref.get("idref")
            href = manifest.get(idref, "")
            if not href.endswith((".html", ".xhtml", ".htm")):
                continue
            full_path = f"{opf_dir}/{href}".lstrip("/")
            try:
                html_bytes = zf.read(full_path)
                parsed = _parse_html(html_bytes)
                if parsed.text.strip():
                    parts.append(parsed.text)
            except Exception:
                continue

    return ParsedDocument(text="\n\n".join(parts))


def _parse_text(content: bytes) -> ParsedDocument:
    text = content.decode("utf-8", errors="replace")
    return ParsedDocument(text=text)
